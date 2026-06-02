"""Generate FPA_Copilot_Presentation.pptx from the Capstone PPT Template."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

_TEMPLATE = os.path.join(os.path.dirname(__file__), "..", "Capstone PPT Template.pptx")
_OUTPUT   = os.path.join(os.path.dirname(__file__), "FPA_Copilot_Presentation.pptx")

# Layout indices confirmed from template inspection
_LAYOUT_TITLE   = 0   # Title Slide  (no placeholders — uses inherited rectangle)
_LAYOUT_CONTENT = 1   # Title and Content (idx 0=title, idx 1=body)
_LAYOUT_SECTION = 8   # Section Header   (idx 0=title, idx 1=body)


def _add_slide(prs, layout_idx):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def _set_title(slide, text):
    slide.shapes.title.text = text


def _set_body(slide, items, placeholder_idx=1):
    """Populate a content placeholder with a list of strings."""
    tf = slide.placeholders[placeholder_idx].text_frame
    tf.clear()
    for i, item in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = item
        para.level = 0


def _add_title_slide(prs, title, subtitle_lines):
    """Title Slide layout has no placeholders — add a centred text box."""
    slide = _add_slide(prs, _LAYOUT_TITLE)
    w = prs.slide_width
    h = prs.slide_height
    # Main title box
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.2), w - Inches(2), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    # Subtitle box
    tb2 = slide.shapes.add_textbox(Inches(1), Inches(4.0), w - Inches(2), Inches(1.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(subtitle_lines):
        para = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        para.text = line
        para.alignment = PP_ALIGN.CENTER
        if para.runs:
            para.runs[0].font.size = Pt(20)
    return slide


# ── Slide content definitions ─────────────────────────────────────────────────

SLIDES = [
    {
        "layout": _LAYOUT_TITLE,
        "title":  "FP&A Narrative Reporting Copilot",
        "body":   ["AI-Powered Variance Analysis & Executive Summaries",
                   "Capstone Project  ·  May 2026"],
    },
    {
        "layout": _LAYOUT_CONTENT,
        "title":  "Agenda",
        "body":   [
            "1.  Problem Statement",
            "2.  Solution Overview",
            "3.  Architecture & 4-Stage Pipeline",
            "4.  Component Breakdown",
            "5.  Prompting Strategy",
            "6.  Streamlit Dashboard Demo",
            "7.  Evaluation Framework & Results",
            "8.  Limitations & Next Steps",
        ],
    },
    {
        "layout": _LAYOUT_CONTENT,
        "title":  "Problem Statement",
        "body":   [
            "Manual variance commentary is slow — analysts spend hours writing what should take minutes",
            "Narratives are inconsistent across months and authors, reducing leadership trust",
            "No built-in fact-checking: LLMs can hallucinate financial figures, risking credibility",
            "Finance teams need a tool that is fast, consistent, and verifiably grounded in the data",
        ],
    },
    {
        "layout": _LAYOUT_CONTENT,
        "title":  "Solution Overview",
        "body":   [
            "Deterministic math via Pandas — every number computed from source Excel, never by the LLM",
            "RAG-grounded narrative via Claude — commentary anchored to finance glossary & policy notes",
            "Built-in Review Mode auditor — a second LLM pass flags unsupported claims & asks questions",
            "Streamlit dashboard — 3-tab UI ready for analyst use in under 30 seconds per report",
        ],
    },
    {
        "layout": _LAYOUT_CONTENT,
        "title":  "Architecture: The WAT Framework",
        "body":   [
            "WORKFLOWS  →  MasterPromt.md  (the SOP defining objective, inputs, constraints)",
            "AGENTS     →  src/engine.py   (4-stage orchestration; LLM decision-making)",
            "TOOLS      →  src/tools.py + src/rag.py  (deterministic math + retrieval)",
            "",
            "Pipeline:  P&L Excel  ──►  Planner  ──►  Compute  ──►  Draft (+RAG)  ──►  Review  ──►  UI",
            "",
            "Key insight: AI handles reasoning; deterministic code handles execution.",
            "5 steps at 90% each = 59% overall if LLM does everything  →  WAT fixes this.",
        ],
    },
    {
        "layout": _LAYOUT_CONTENT,
        "title":  "Component Breakdown",
        "body":   [
            "src/tools.py   — Pandas variance engine; unified actual−target formula for all categories",
            "src/rag.py     — ChromaDB PersistentClient; indexes finance glossary + policy notes",
            "src/engine.py  — 4-stage pipeline (Planner → Compute → Draft → Review); 2 LLM calls",
            "app.py         — Streamlit dashboard; st.session_state cache prevents redundant API calls",
            "eval.py        — LLM-as-judge evaluation; coverage + accuracy scores (1–5 each)",
            "tests/         — 7 pytest cases; sign-convention & variance math verification",
        ],
    },
    {
        "layout": _LAYOUT_CONTENT,
        "title":  "Prompting Strategy",
        "body":   [
            "DRAFT STAGE — System role: 'FP&A analyst writing for senior leadership'",
            "  • RAG context injected into system message before the task begins",
            "  • Explicit rules: quantify in INR lakhs + %, top 2–3 drivers, no invented figures",
            "",
            "REVIEW STAGE — System role: 'Strict financial audit reviewer'",
            "  • Verified driver data passed as ground truth",
            "  • Outputs structured bullet risks + clarifying questions; flags 'Data Insufficient'",
            "",
            "WHY TWO CALLS: One prompt asking Claude to both write and critique produces leniency.",
            "Separate personas create genuine adversarial tension — drafter vs. skeptic.",
        ],
    },
    {
        "layout": _LAYOUT_CONTENT,
        "title":  "Evaluation Framework & Results",
        "body":   [
            "3 test prompts from evaluation_prompts.xlsx:",
            "  E01 — Nov 2025 vs Budget   (expected: S&M overspend + revenue recovery)",
            "  E02 — Sep 2025 vs Forecast  (expected: Product Revenue dip + COGS pressure)",
            "  E03 — Oct 2025 management  (expected: cautious language, no market claims)",
            "",
            "LLM-as-judge scores two dimensions (1–5 each):",
            "  Coverage  — does narrative capture the expected focus drivers?",
            "  Accuracy  — are all figures verifiable from source data?",
            "",
            "Limitation: judge uses same model (self-evaluation bias) → use Haiku for production judge.",
        ],
    },
    {
        "layout": _LAYOUT_CONTENT,
        "title":  "Limitations & Next Steps",
        "body":   [
            "LIMITATIONS:",
            "  • Synthetic data only — not validated on real enterprise P&L edge cases",
            "  • Self-evaluation bias in LLM judge — scores may be inflated",
            "  • No multi-month trend analysis or YTD context",
            "  • No user authentication on Streamlit app",
            "",
            "NEXT STEPS:",
            "  • Add MoM/YTD trend charts in Tab 1 using st.line_chart",
            "  • One-click export narrative to Word/PDF for distribution",
            "  • Replace judge model with claude-haiku-4-5 for independent scoring",
            "  • GitHub Actions CI/CD running pytest on every push",
        ],
    },
    {
        "layout": _LAYOUT_SECTION,
        "title":  "Thank You & Q&A",
        "body":   [
            "Project: FP&A Narrative Reporting Copilot",
            "Stack: Python · Streamlit · ChromaDB · Claude (claude-sonnet-4-6)",
            "Repo: /AIAgents/BIAproject/fpa_copilot/",
        ],
    },
]


def build_presentation():
    prs = Presentation(_TEMPLATE)

    # Remove all existing placeholder slides from the template.
    # _sldIdLst holds <p:sldId> elements (not the slide XML itself).
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        sldIdLst.remove(sldId)

    # Add slides from content definitions
    for slide_def in SLIDES:
        layout_idx = slide_def["layout"]
        body_items = slide_def.get("body", [])

        if layout_idx == _LAYOUT_TITLE:
            # Title Slide has no placeholders — use text box approach
            _add_title_slide(prs, slide_def["title"], body_items)
            continue

        slide = _add_slide(prs, layout_idx)

        try:
            _set_title(slide, slide_def["title"])
        except Exception:
            pass

        if body_items:
            try:
                _set_body(slide, body_items, placeholder_idx=1)
            except (KeyError, IndexError):
                pass

    prs.save(_OUTPUT)
    print(f"Saved: {_OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
